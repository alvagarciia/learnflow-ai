"""
File processing service for extracting text from various document formats.
Handles PDF, DOCX, PPTX extraction with memory-safe streaming.
"""

import io
from typing import Dict, List, Optional
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


class FileProcessor:
    """Extract text content from uploaded files"""
    
    # Maximum file sizes (in bytes)
    MAX_PDF_SIZE = 15 * 1024 * 1024  # 15MB
    MAX_DOCX_SIZE = 15 * 1024 * 1024  # 15MB
    MAX_PPTX_SIZE = 30 * 1024 * 1024  # 30MB
    
    @staticmethod
    def extract_from_pdf(file_stream: io.BytesIO, filename: str) -> Dict[str, any]:
        """
        Extract text from PDF file.
        
        Args:
            file_stream: Binary stream of PDF file
            filename: Original filename for metadata
            
        Returns:
            Dict with extracted text and metadata
            
        Raises:
            ValueError: If PDF is corrupted or unreadable
        """
        try:
            reader = PdfReader(file_stream)
            
            if reader.is_encrypted:
                raise ValueError(f"PDF '{filename}' is encrypted and cannot be processed")
            
            text_content = []
            total_pages = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, start=1):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(page_text)
                except Exception as e:
                    # Log but continue processing other pages
                    print(f"Warning: Could not extract page {page_num} from {filename}: {str(e)}")
                    continue
            
            if not text_content:
                raise ValueError(f"PDF '{filename}' contains no extractable text")
            
            return {
                "source": filename,
                "type": "pdf",
                "text": "\n\n".join(text_content),
                "metadata": {
                    "pages": total_pages,
                    "extracted_pages": len(text_content)
                }
            }
            
        except Exception as e:
            raise ValueError(f"Failed to process PDF '{filename}': {str(e)}")
    
    @staticmethod
    def extract_from_docx(file_stream: io.BytesIO, filename: str) -> Dict[str, any]:
        """
        Extract text from DOCX file.
        
        Args:
            file_stream: Binary stream of DOCX file
            filename: Original filename for metadata
            
        Returns:
            Dict with extracted text and metadata
            
        Raises:
            ValueError: If DOCX is corrupted or unreadable
        """
        try:
            doc = Document(file_stream)
            
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            
            if not paragraphs:
                raise ValueError(f"DOCX '{filename}' contains no extractable text")
            
            return {
                "source": filename,
                "type": "docx",
                "text": "\n\n".join(paragraphs),
                "metadata": {
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables)
                }
            }
            
        except Exception as e:
            raise ValueError(f"Failed to process DOCX '{filename}': {str(e)}")
    
    @staticmethod
    def extract_from_pptx(file_stream: io.BytesIO, filename: str) -> Dict[str, any]:
        """
        Extract text from PPTX file.
        Extracts slide titles, body text, and notes.
        
        Args:
            file_stream: Binary stream of PPTX file
            filename: Original filename for metadata
            
        Returns:
            Dict with extracted text and metadata
            
        Raises:
            ValueError: If PPTX is corrupted or unreadable
        """
        try:
            prs = Presentation(file_stream)
            
            slide_contents = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                # Extract title
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_text.append(f"Slide {slide_num}: {title}")
                
                # Extract body text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if text and (not slide.shapes.title or text != slide.shapes.title.text):
                            slide_text.append(text)
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"Notes: {notes_text}")
                
                if slide_text:
                    slide_contents.append("\n".join(slide_text))
            
            if not slide_contents:
                raise ValueError(f"PPTX '{filename}' contains no extractable text")
            
            return {
                "source": filename,
                "type": "pptx",
                "text": "\n\n---\n\n".join(slide_contents),
                "metadata": {
                    "total_slides": len(prs.slides),
                    "extracted_slides": len(slide_contents)
                }
            }
            
        except Exception as e:
            raise ValueError(f"Failed to process PPTX '{filename}': {str(e)}")
    
    @classmethod
    def process_file(cls, file_stream: io.BytesIO, filename: str, file_size: int) -> Dict[str, any]:
        """
        Main entry point for file processing.
        Routes to appropriate extractor based on file extension.
        
        Args:
            file_stream: Binary stream of file
            filename: Original filename
            file_size: File size in bytes for validation
            
        Returns:
            Dict with extracted text and metadata
            
        Raises:
            ValueError: If file type unsupported or validation fails
        """
        extension = filename.lower().split('.')[-1]
        
        # Validate file type and size
        if extension == 'pdf':
            if file_size > cls.MAX_PDF_SIZE:
                raise ValueError(f"PDF file too large. Maximum size: {cls.MAX_PDF_SIZE / (1024*1024)}MB")
            return cls.extract_from_pdf(file_stream, filename)
        
        elif extension in ['doc', 'docx']:
            if file_size > cls.MAX_DOCX_SIZE:
                raise ValueError(f"DOCX file too large. Maximum size: {cls.MAX_DOCX_SIZE / (1024*1024)}MB")
            return cls.extract_from_docx(file_stream, filename)
        
        elif extension in ['ppt', 'pptx']:
            if file_size > cls.MAX_PPTX_SIZE:
                raise ValueError(f"PPTX file too large. Maximum size: {cls.MAX_PPTX_SIZE / (1024*1024)}MB")
            return cls.extract_from_pptx(file_stream, filename)
        
        else:
            raise ValueError(f"Unsupported file type: .{extension}")